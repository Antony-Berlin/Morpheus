#helper function here
import PyPDF2
from pptx import Presentation
from dotenv import load_dotenv

import google.generativeai as genai
import os, json, time
import csv


def init():
    load_dotenv()
    global model 

    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
    model = genai.GenerativeModel("gemini-1.5-flash")
    print("generative AI client initialized ...")

def chat(system_message : str , prompt : str):

    response = None
    for _ in range(3):  # Retry up to 3 times
        try:
            response = model.generate_content(system_message + "\n" + prompt)
            break  # Exit loop if successful
        except Exception as e:
            print(f"Error occurred: {e}. Retrying in 10 seconds...")
            time.sleep(15)

    return response.text

def extract_text(file_name):
    print("extracting text from file...")
    file_extension = os.path.splitext(file_name)[1].lower()
    current_dir = os.path.dirname(__file__)
    file_path = os.path.join(current_dir, '..', 'Documents', file_name)
    
    if file_extension == '.ppt' or file_extension == '.pptx':
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                    text_runs.append("-"*10)
        text =  "\n".join(text_runs)
    
    elif file_extension == '.pdf':
        text_runs = []
        with open(file_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text_runs.append(page.extract_text())
                text_runs.append("-"*10)
        text = "\n".join(text_runs)
    
    else:
        raise ValueError("Unsupported file format. Please provide a PPT, PPTX, or PDF file.")
        text = text.replace("\n", " ")
    return text


def get_topics(text, no_of_questions):
    print("generating topics...")
    no_of_questions = int(no_of_questions)
    if no_of_questions == 0:
        return []
    if(no_of_questions%2 ==0):
        max = no_of_questions/2
    elif(no_of_questions%3 == 0):
        max = no_of_questions/3
    else:
        max = no_of_questions/5
    print(max)
    system_message = "you are expert in providing the list of topics given a paragraph.you're response is always comma seperated topics no other sentence is required. help the student by giving the topics from his notes"
    prompt = "give me only "+ str(max) + " major topics from the below student notes:\n" + text
    response = chat(system_message=system_message, prompt=prompt)
    
    topics = response.split(",")
    print(topics)
    print(len(topics))
    return topics[:int(max)]



def get_questions_and_answers(topic, text, about_proff, proff_sample_question, no_of_questions, prev_questions):
  
    # request_count = 0
    # max_requests_per_minute = 15
    # start_time = time.time()
    qa = []

    
    # if request_count >= max_requests_per_minute:
    #     elapsed_time = time.time() - start_time
    #     remaining_time = 65 - elapsed_time
    #     if remaining_time > 0:
    #         print(f"Rate limit reached. Waiting for {remaining_time:.2f} seconds...")
    #         time.sleep(remaining_time)
    #     start_time = time.time()  # Reset the start time after waiting
    #     request_count = 0  # Reset the request count after waiting


    System_message = "you are expert in mimicing the professor in generating question, options and its answer.Provide the JSON response directly without wrapping it in backticks or marking it as a code block. "

    prompt = (
        
        "example: [{\"question\": \"what is the capital of india?\", \"options\": [\"delhi\", \"mumbai\", \"kolkata\", \"chennai\"], \"answer\": \"delhi\"}, {\"question\": \"what is the capital of usa?\", \"options\": [\"new york\", \"washington dc\", \"los angeles\", \"chicago\"], \"answer\": \"washington dc\"}]"+
        "content: " + text + "\n" +
        "topics: " + topic + "\n" +
        "about professor: " + about_proff + "\n" +
        "professor sample question: " + proff_sample_question + "\n" +
        "no of questions: " + str(no_of_questions) + "\n" +
        "previous questions: " + str(prev_questions) + "\n" +
        "as a professor u know the content and now return the list of mcq questions with options and answers as dictionary object, the response must only have the array, shouldn't have any context or extra sentence. avoid any previous questions repetation \n"

    )

    response = chat(system_message=System_message, prompt=prompt)
    
    try:
        response_data = json.loads(response)
        qa.extend(response_data)
    except json.JSONDecodeError as e:
        print(f"Failed to decode JSON response: {e}")
        print(f"Response text: {response}")
        return {"error": "Failed to decode JSON response"}, 400
    except Exception as e:
        print(f"An error occurred: {e}")
        print(f"Response text: {response}")
        return {"error": "An error occurred"}, 400

    return qa[:(int(no_of_questions)+1)], 200

def handle_file_upload(file):
    if file.filename == '':
        return {"error": "No selected file"}, 400
    if file and allowed_file(file.filename):
        filename = file.filename
        file_path = os.path.join(os.getcwd(), "Documents", filename)
        
        if os.path.exists(file_path):
            return {"error": "File already exists"}, 409
        
        file.save(file_path)
        print("Document uploaded successfully")
        return {"message": "File uploaded successfully", "file_name": filename}, 201
    else:
        return {"error": "File type not allowed"}, 415

def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def add_proff_profile(prof_name, about_prof, prof_sample_question):
    file_path = os.path.join(os.getcwd(), "proff_profile.csv")
    file_exists = os.path.isfile(file_path)

    # Ensure the sample question is a single string
    prof_sample_question = prof_sample_question.replace("\n", " ").replace("\r", " ")

    if file_exists:
        with open(file_path, mode='r', newline='') as file:
            reader = csv.reader(file)
            for row in reader:
                if row and row[0] == prof_name:
                    return {"error": "Professor profile already exists"}, 409

    with open(file_path, mode='a', newline='') as file:
        writer = csv.writer(file)
        if not file_exists:
            writer.writerow(["prof_name", "about_prof", "prof_sample_question"])
        writer.writerow([prof_name, about_prof, prof_sample_question])

    return {"message": "Professor profile added successfully", "proff_name": prof_name}, 201

def get_proff_profile_data(prof_name):
    file_path = os.path.join(os.getcwd(), "proff_profile.csv")
    if not os.path.exists(file_path):
        return {"error": "Professor profile does not exist"}, 404

    with open(file_path, mode='r', newline='') as file:
        reader = csv.reader(file)
        for row in reader:
            if row and row[0] == prof_name:
                return {"prof_name": row[0], "about_prof": row[1], "prof_sample_question": row[2]}, 200

    return {"error": "Professor profile does not exist"}, 404

def delete_proff_profile_data(prof_name):
    file_path = os.path.join(os.getcwd(), "proff_profile.csv")
    if not os.path.exists(file_path):
        return {"error": "Professor profile does not exist"}, 404

    with open(file_path, mode='r', newline='') as file:
        reader = csv.reader(file)
        rows = list(reader)

    if len(rows) <= 1:
        return {"error": "Professor profile does not exist"}, 404

    with open(file_path, mode='w', newline='') as file:
        writer = csv.writer(file)
        for row in rows:
            if row and row[0] != prof_name:
                writer.writerow(row)
    return {"message": "Professor profile deleted successfully"}, 200

def get_all_proff_name():
    file_path = os.path.join(os.getcwd(), "proff_profile.csv")
    if not os.path.exists(file_path):
        return {"error": "Professor profile does not exist"}, 404

    with open(file_path, mode='r', newline='') as file:
        reader = csv.reader(file)
        next(reader)
        return {"professors": [row[0] for row in reader]},200

    return {"error": "Professor profile does not exist"}, 404

